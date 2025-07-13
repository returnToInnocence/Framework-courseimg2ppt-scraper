from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from PIL import Image
import time
from natsort import natsorted
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from selenium.webdriver.common.action_chains import ActionChains
import os
import shutil
import json
from configparser import ConfigParser

def load_config():
    """加载配置文件"""
    config = ConfigParser()
    config.read('config.ini')
    return config

def clear_folder(folder_path):
    if os.path.exists(folder_path):
        shutil.rmtree(folder_path)
    os.makedirs(folder_path)

def take_screenshot(driver, save_path, width, height, crop_area=None):
    driver.set_window_size(width, height)
    time.sleep(2)
    driver.save_screenshot(save_path)
    time.sleep(2)
    if crop_area:
        img = Image.open(save_path)
        cropped_img = img.crop(crop_area)
        cropped_img.save(save_path)
        time.sleep(2)

def create_ppt_with_images(image_folder, ppt_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    image_files = natsorted([f for f in os.listdir(image_folder) if f.endswith('.png')])

    for img_path in image_files:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        img_full_path = os.path.join(image_folder, img_path)

        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_full_path, left, top, width=prs.slide_width, height=prs.slide_height)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    prs.save(ppt_path)

def main():
    # 加载配置
    config = load_config()
    
    # 配置Chrome选项
    chrome_options = Options()
    
    # 从配置文件读取Chrome路径
    chrome_binary_path = config.get('chrome', 'binary_location')
    chromedriver_path = config.get('chrome', 'chromedriver_path')
    
    chrome_options.binary_location = chrome_binary_path
    chrome_options.add_argument('--no-sandbox')
    chrome_options.add_argument('--disable-dev-shm-usage')
    chrome_options.add_argument('--disable-gpu')
    
    # 创建Service对象
    service = Service(executable_path=chromedriver_path)
    
    print("正在启动Chrome浏览器...")
    
    # 初始化Chrome浏览器
    browser = webdriver.Chrome(service=service, options=chrome_options)
    browser.maximize_window()
    
    print("Chrome浏览器启动成功，正在访问网站...")
    
    time.sleep(1)

    # 从配置文件读取URL
    url = config.get('app', 'target_url')
    browser.get(url)
    
    print("网站访问成功，正在设置页面缩放...")

    browser.execute_script("document.body.style.zoom='80%'")
    time.sleep(1)
    
    # 登录逻辑
    username_xpath = config.get('xpaths', 'username_xpath')
    password_xpath = config.get('xpaths', 'password_xpath')
    login_button_xpath = config.get('xpaths', 'login_button_xpath')
    
    username_path = browser.find_element(By.XPATH, username_xpath)
    password_path = browser.find_element(By.XPATH, password_xpath)
    login_in = browser.find_element(By.XPATH, login_button_xpath)
    
    # 从配置文件读取用户名和密码
    username = config.get('credentials', 'username')
    password = config.get('credentials', 'password')
    
    username_path.send_keys(username)
    time.sleep(1)
    password_path.send_keys(password)
    login_in.click()
    time.sleep(4)

    # 系统导航逻辑
    jiaopei_system_xpath = config.get('xpaths', 'jiaopei_system_xpath')
    jiaopei_system = browser.find_element(By.XPATH, jiaopei_system_xpath)
    jiaopei_system.click()
    time.sleep(1)
    
    jiaoyan_button_xpath = config.get('xpaths', 'jiaoyan_button_xpath')
    jiaoyan_button = browser.find_element(By.XPATH, jiaoyan_button_xpath)
    actions = ActionChains(browser)
    actions.move_to_element(jiaoyan_button).perform()
    time.sleep(1)
    
    course_management_xpath = config.get('xpaths', 'course_management_xpath')
    course_management_button = browser.find_element(By.XPATH, course_management_xpath)
    course_management_button.click()
    time.sleep(1)
    
    # 搜索课程
    course_input_xpath = config.get('xpaths', 'course_input_xpath')
    search_button_xpath = config.get('xpaths', 'search_button_xpath')
    
    course_input = browser.find_element(By.XPATH, course_input_xpath)
    course_input.send_keys(config.get('app', 'search_keyword'))
    search = browser.find_element(By.XPATH, search_button_xpath)
    search.click()
    time.sleep(2)
    
    # 课程选择
    course_xpath = config.get('xpaths', 'course_xpath')
    course355 = browser.find_element(By.XPATH, course_xpath)
    course355.click()
    time.sleep(1)
    
    classprogress_xpath = config.get('xpaths', 'classprogress_xpath')
    classprogress = browser.find_element(By.XPATH, classprogress_xpath)
    classprogress.click()
    time.sleep(2)

    # 从配置文件读取路径设置
    width = config.getint('screenshot', 'width')
    height = config.getint('screenshot', 'height')
    image_folder = config.get('paths', 'image_folder')
    output_folder = config.get('paths', 'output_folder')

    m = 1  # 课程单元数
    counts = 1  # 课程数量
    
    # 其余逻辑保持不变...
    while m <= 1:
        if m != 1:
            unit_xpath = config.get('xpaths', 'unit_xpath').format(m=m)
            unit = browser.find_element(By.XPATH, unit_xpath)
            time.sleep(1)
            unit.click()
            time.sleep(2)

        j = 1
        while j <= 10:
            i = 1
            clear_folder(image_folder)
            
            course_item_xpath = config.get('xpaths', 'course_item_xpath').format(m=m, j=j)
            course_id = browser.find_element(By.XPATH, course_item_xpath)
            course_id.click()
            time.sleep(2)
            
            try:
                element_xpath = config.get('xpaths', 'element_xpath').format(m=m, j=j)
                element = WebDriverWait(browser, 20).until(
                    EC.presence_of_element_located((By.XPATH, element_xpath))
                )

                element_text = element.text
                print(f"元素文本内容: {element_text}")

                main_window = browser.current_window_handle

                course_ppt_xpath = config.get('xpaths', 'course_ppt_xpath')
                course_ppt = browser.find_element(By.XPATH, course_ppt_xpath)
                course_ppt.click()
                time.sleep(5)

                ppt_windows = browser.window_handles
                browser.switch_to.window(ppt_windows[-1])

                browser.execute_script("document.body.style.zoom='80%'")
                pc = config.getint('app', 'max_pages')
                
                crop_area_config = config.get('screenshot', 'crop_area').split(',')
                crop_area = tuple(map(int, crop_area_config))
                
                while i < pc:
                    print(i)
                    screenshot_path = os.path.join(image_folder, f'{i}screenshot.png')
                    take_screenshot(browser, screenshot_path, width, height, crop_area)
                    
                    next_page_xpath = config.get('xpaths', 'next_page_xpath')
                    next_page = browser.find_element(By.XPATH, next_page_xpath)
                    next_page.click()
                    time.sleep(2)
                    i += 1
                
                ppt_path = os.path.join(output_folder, f'课程{counts} {element_text}.pptx')
                create_ppt_with_images(image_folder, ppt_path)
                j += 1
                counts += 1
                browser.switch_to.window(main_window)
            except:
                print(f"元素未找到")
                break
        m += 1

if __name__ == "__main__":
    main()